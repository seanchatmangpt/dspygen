import dspy
from typing import List, Union, Optional

import pandas as pd
import pm4py
from pydantic import BaseModel

LogType = Union[pd.DataFrame, pm4py.objects.log.obj.EventLog, pm4py.objects.log.obj.EventStream]


class Slide(DSLModel):
    title: str
    content: Union[str, List[str]]
    image_path: Optional[str] = None


class Presentation(DSLModel):
    title: str
    subtitle: str
    slides: List[Slide]
    file_path: Optional[str] = None

    def insert_slide_at(self, slide: Slide, position: int) -> None:
        self.slides.insert(position, slide)

    def to_ppt(self, file_path: str = None):
        if file_path:
            self.file_path = file_path

        # install python-pptx
        from pptx import Presentation as PPTPresentation
        from pptx.util import Inches
        from PIL import Image

        prs = PPTPresentation()

        # Title Slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = self.title
        subtitle.text = self.subtitle

        # Add Slides
        for slide_data in self.slides:
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            content = slide.placeholders[1]
            title.text = slide_data.title

            if isinstance(slide_data.content, list):
                content.text = "\n".join(slide_data.content)
            else:
                content.text = slide_data.content

            # Insert image if provided
            if slide_data.image_path:
                img = Image.open(slide_data.image_path)
                width, height = img.size
                aspect_ratio = width / height

                max_width = Inches(6)  # Maximum width for the image
                max_height = Inches(4.5)  # Maximum height for the image

                # Calculate the dimensions to fit within the slide while maintaining aspect ratio
                if aspect_ratio > 1:  # Landscape
                    if width > max_width:
                        width = max_width
                        height = width / aspect_ratio
                else:  # Portrait
                    if height > max_height:
                        height = max_height
                        width = height * aspect_ratio

                left = (prs.slide_width - width) / 2  # Center the image horizontally
                top = (prs.slide_height - height) / 2  # Center the image vertically

                slide.shapes.add_picture(slide_data.image_path, left, top, width, height)

        prs.save(file_path)
        print(f"Presentation saved as '{file_path}'")


class GenerateGlossarySlide(dspy.Signature):
    """
    Generates a glossary slide for the provided presentation JSON string.
    The format should be:

    - {TERM}: {DEFINITION}
    """
    presentation_json_str = dspy.InputField(desc="Presentation JSON string.")
    glossary_slide = dspy.OutputField(desc="Generated glossary slide content. NO FORMATING OR MARKDOWN. PLAINTEXT ONLY", prefix="```plaintext\nGlossary Slide:\n\n")


class GlossaryModule(dspy.Module):
    """GlossaryModule"""

    def __init__(self, **forward_args):
        super().__init__()
        self.forward_args = forward_args
        self.output = None

    def forward(self, presentation_json_str):
        pred = dspy.Predict(GenerateGlossarySlide)
        self.output = pred(presentation_json_str=presentation_json_str).glossary_slide
        return self.output


def glossary_call(presentation: Presentation) -> Slide:
    presentation_json_str = presentation.json()
    glossary = GlossaryModule()
    glossary_slide_content = glossary.forward(presentation_json_str=presentation_json_str)
    return Slide(title="Glossary of Terms", content=glossary_slide_content)


def main():
    """Main function"""
    from dspygen.utils.dspy_tools import init_ol
    init_ol()

    # Create an example presentation
    example_presentation = Presentation(
        title="Analysis of Car Insurance Claims Event Log",
        slides=[
            Slide(title="Introduction", content="Overview of the project and dataset."),
            Slide(title="Directly-Follows Graph (DFG) Analysis",
                  content="Visual representation and analysis of the DFG."),
            Slide(title="Variants Analysis", content="Analysis of the most common sequences of actions."),
            Slide(title="Petri Net Model of the Process", content="Visual representation of the discovered Petri net."),
            Slide(title="Event Log Attributes", content="Insights into the data characteristics."),
            Slide(title="Conclusion", content=[
                "1. DFG provides an overview of the process flow.",
                "2. Variants analysis reveals the most common sequences of actions.",
                "3. Petri net model highlights the detailed process structure.",
                "4. Event log attributes give insights into the data characteristics."
            ])
        ]
    )

    # Generate the glossary slide
    glossary_slide = glossary_call(example_presentation)

    # Insert the glossary slide into the presentation
    example_presentation.insert_slide_at(glossary_slide, position=1)

    # Print the updated presentation structure
    for slide in example_presentation.slides:
        print(f"Title: {slide.title}")
        print(f"Content: {slide.content}")
        print()

    # Create the PowerPoint presentation with the glossary slide
    example_presentation.to_ppt(file_path="car_insurance_claims_analysis_with_glossary.pptx")


if __name__ == '__main__':
    main()
